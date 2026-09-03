"""v1.6.0 — File read/write/edit tools.

Three tools the agent can call for file operations:

  read_file   — read a file, return content as text
  write_file  — write content to a file (overwrite)
  edit_file   — find-and-replace within a file

All three enforce a **working-directory allowlist** — the agent can
only touch files under an allowed root. This prevents the agent from
reading ``~/.ssh/id_rsa`` or writing to ``/etc/passwd``.

Design (Qian control theory):
  - 可控性: every read/write is scoped to allowed_dirs
  - 稳定性: size caps on read/write prevent memory exhaustion
  - 早纠偏: write_file refuses to create directories outside allowlist
  - 层次化: these tools compose with the sandbox (subprocess) layer
"""
from __future__ import annotations

import logging
import os
import re
from pathlib import Path
from typing import Any, Sequence

from .registry import Tool

logger = logging.getLogger(__name__)

_MAX_READ_BYTES = 500_000   # 500 KB read cap
_MAX_WRITE_BYTES = 500_000  # 500 KB write cap


def _resolve_in_allowlist(
    path: str | Path,
    allowed_dirs: Sequence[str | Path],
) -> Path:
    """Resolve ``path`` and verify it's inside one of ``allowed_dirs``.

    Relative paths are resolved against the **first** allowed dir (the
    active workspace), so models can write ``out.md`` instead of a full
    absolute path and still land in the user's project.

    Hallucinated-username auto-correction: models frequently write to
    ``/Users/xiaoming/…`` (the canonical Chinese textbook username) or
    other nonexistent users under /Users. When the first path segment
    under /Users names a user that does not exist on this machine, the
    path is transparently rewritten to the REAL home and resolution
    retried — the tool result tells the model the corrected path so
    later references stay consistent.

    Raises ``PermissionError`` if still outside. Denied attempts are
    logged so a bypass attempt is auditable even when the caller only
    returns ``{"error": ...}`` to the model.
    """
    raw = Path(path).expanduser()

    # Hallucinated home-dir correction (before resolution — the
    # nonexistent segment can't be resolved faithfully anyway). Covers
    # the three corpus-driven patterns: /Users/<name>, /home/<name>,
    # C:\Users\<name> — xiaoming is THE Chinese textbook placeholder.
    real_home = Path.home()
    try:
        parts = raw.parts
        if (
            len(parts) >= 3 and parts[0] == "/" and parts[1] in ("Users", "home")
            and parts[2] != real_home.name
            and not (Path("/") / parts[1] / parts[2]).exists()
        ):
            raw = real_home.joinpath(*parts[3:])
        else:
            rs = str(raw)
            for win_prefix in ("C:\\Users\\", "C:/Users/"):
                if rs.startswith(win_prefix):
                    segs = re.split(r"[\\/]", rs)
                    if (
                        len(segs) >= 3 and segs[2] != real_home.name
                        and not Path("C:/Users", segs[2]).exists()
                    ):
                        raw = real_home.joinpath(*segs[3:])
                    break
    except Exception:
        pass

    if not raw.is_absolute():
        if allowed_dirs:
            base = Path(allowed_dirs[0]).expanduser().resolve()
            p = (base / raw).resolve()
        else:
            p = raw.resolve()
    else:
        p = raw.resolve()
    for allowed in allowed_dirs:
        a = Path(allowed).expanduser().resolve()
        if p == a or a in p.parents:
            return p
    allowed_list = [str(Path(d).expanduser().resolve()) for d in allowed_dirs]
    logger.warning(
        "allowlist denial: path=%s allowed_dirs=%s",
        p,
        allowed_list,
    )
    raise PermissionError(
        f"Path '{p}' is outside allowed directories: {allowed_list}"
    )


def _extract_docx_text(raw_bytes: bytes) -> str:
    """Extract text from a .docx (Office Open XML) document via python-docx.

    Returns the extracted text (paragraphs + tables), or an empty string on
    failure so the caller can decide how to report it.
    """
    try:
        import io as _io
        from docx import Document as _DocxDocument
        doc = _DocxDocument(_io.BytesIO(raw_bytes))
        parts: list[str] = []
        for para in doc.paragraphs:
            if para.text and para.text.strip():
                parts.append(para.text)
        for ti, table in enumerate(doc.tables):
            parts.append(f"\n[Table {ti + 1}]")
            for row in table.rows:
                parts.append(" | ".join(c.text.strip() for c in row.cells))
        return "\n".join(parts).strip()
    except Exception:
        return ""


# --------------------------------------------------------------------------- #
# ReadFileTool
# --------------------------------------------------------------------------- #


class ReadFileTool(Tool):
    """Read a text file and return its content."""

    name = "read_file"
    description = (
        "Read a text file and return its content. "
        "The path must be inside an allowed working directory."
    )

    def __init__(self, allowed_dirs: Sequence[str | Path] | None = None) -> None:
        self._allowed_dirs = list(allowed_dirs or [os.getcwd()])

    @property
    def parameters_schema(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Path to the file to read.",
                },
                "offset": {
                    "type": "integer",
                    "description": "Line number to start from (1-indexed). Default 1.",
                },
                "limit": {
                    "type": "integer",
                    "description": "Max lines to read. Default 500.",
                },
            },
            "required": ["path"],
        }

    def _read_attachment(self, att: dict[str, Any]) -> dict[str, Any]:
        """Decode an inline attachment (from chat composer) and return its
        content as a string. Supports text files (utf-8) and a metadata
        header for binary files like images/PDFs."""
        name = att.get("name", "attachment")
        mime = att.get("mimeType", "")
        data = att.get("data", "")
        if not data:
            return {"error": f"no data for attachment {name}"}
        # data may be a data URL (data:<mime>;base64,<...>) — keep the
        # metadata but truncate the body for very large files.
        if data.startswith("data:") and "," in data:
            header, body = data.split(",", 1)
            # If it's a text-y mime, try to decode and return text.
            if mime.startswith("text/") or mime in ("application/json", "application/xml"):
                import base64 as _b64
                try:
                    return {"path": att.get("id") or name,
                            "content": _b64.b64decode(body).decode("utf-8", errors="replace")}
                except Exception as e:
                    return {"error": f"failed to decode {name}: {e}"}
            # PDF: extract real text using pypdf
            if mime == "application/pdf" or name.lower().endswith(".pdf"):
                try:
                    import base64 as _b64
                    import io as _io
                    from pypdf import PdfReader as _PdfReader
                    raw_bytes = _b64.b64decode(body)
                    reader = _PdfReader(_io.BytesIO(raw_bytes))
                    pages_text = []
                    for i, page in enumerate(reader.pages):
                        try:
                            t = page.extract_text() or ""
                        except Exception:
                            t = ""
                        pages_text.append(f"--- Page {i + 1} ---\n{t}")
                    text = "\n\n".join(pages_text).strip()
                    if text:
                        # Truncate very large PDFs to keep response manageable
                        max_chars = 30_000
                        if len(text) > max_chars:
                            text = text[:max_chars] + f"\n\n[truncated at {max_chars} chars of {len(text)} total]"
                        return {
                            "path": att.get("id") or name,
                            "content": text,
                        }
                    return {
                        "path": att.get("id") or name,
                        "content": f"[PDF text extraction returned no content: {name} (scanned/image-only PDF?)]"
                    }
                except Exception as e:
                    return {
                        "path": att.get("id") or name,
                        "content": f"[PDF parse error: {e}]"
                    }
            # Binary: return metadata; LLM can decide what to do.
            # Excel xlsx — extract as text via openpyxl
            if name.lower().endswith(".xlsx") or name.lower().endswith(".xls") or mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                import base64 as _b64
                raw_bytes = _b64.b64decode(body)
                try:
                    import io as _io, openpyxl as _xl
                    wb = _xl.load_workbook(_io.BytesIO(raw_bytes), data_only=True, read_only=True)
                    parts = []
                    for sheet_name in wb.sheetnames:
                        ws = wb[sheet_name]
                        rows = list(ws.iter_rows(values_only=True))
                        if not rows: continue
                        parts.append(f"## Sheet: {sheet_name} ({len(rows)} rows, {len(rows[0])} cols)")
                        header = " | ".join(str(c or "") for c in rows[0])
                        sep = " | ".join("---" for _ in rows[0])
                        body_rows = []
                        for row in rows[1:]:
                            body_rows.append(" | ".join(str(c or "") for c in row))
                        parts.append(f"| {header} |\n| {sep} |\n" + "\n".join(f"| {r} |" for r in body_rows))
                    wb.close()
                    content = "\n\n".join(parts)
                    return {"path": att.get("id") or name, "content": content[:60_000] or "[empty xlsx]"}
                except Exception as _xe:
                    return {"path": att.get("id") or name, "content": f"[xlsx parse error: {_xe}]"}
            # Word .docx — extract text via python-docx
            if name.lower().endswith(".docx") or mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                import base64 as _b64
                raw_bytes = _b64.b64decode(body)
                text = _extract_docx_text(raw_bytes)
                if text:
                    max_chars = 30_000
                    if len(text) > max_chars:
                        text = text[:max_chars] + f"\n\n[truncated at {max_chars} chars of {len(text)} total]"
                    return {"path": att.get("id") or name, "content": text}
                return {
                    "path": att.get("id") or name,
                    "content": f"[docx parse returned no text: {name} (image-only or corrupted?)]",
                }
            if name.lower().endswith(".doc"):
                return {
                    "path": att.get("id") or name,
                    "content": f"[.doc (legacy Word) not supported: {name} — please convert to .docx]",
                }
            return {
                "path": att.get("id") or name,
                "content": f"[binary file: {name}, type: {mime}, size: {len(body)} base64 chars — describe what you see or do not try render]",
            }
        # Raw text fallback.
        return {"path": att.get("id") or name, "content": data}

    def __call__(self, **kwargs: Any) -> dict[str, Any]:
        path_str = kwargs.get("path", "")
        if not path_str:
            return {"error": "missing 'path'"}

        # Virtual path scheme for inline attachments sent from the chat
        # composer. We can't use the real disk path because the user's
        # file API may not have a path in Electron. Instead we look up
        # the attachment in a module-level registry keyed by id.
        if path_str.startswith("attachment://"):
            from . import inline_attachments
            att_id = path_str[len("attachment://"):]
            att = inline_attachments.get(att_id)
            if not att:
                return {"error": f"attachment not found: {att_id}"}
            return self._read_attachment(att)

        try:
            p = _resolve_in_allowlist(path_str, self._allowed_dirs)
        except PermissionError as e:
            logger.info("read_file denied: %s", e)
            return {"error": str(e)}

        if not p.exists():
            return {"error": f"file not found: {p}"}
        if not p.is_file():
            return {"error": f"not a file: {p}"}

        offset = max(1, int(kwargs.get("offset", 1)))
        limit = min(2000, max(1, int(kwargs.get("limit", 500))))

        # Detect PDF by extension and use pypdf to extract text instead of
        # reading raw binary bytes as UTF-8 (which produces garbage).
        if str(p).lower().endswith(".pdf"):
            try:
                from pypdf import PdfReader as _PdfReader
                reader = _PdfReader(str(p))
                pages_text = []
                for i, page in enumerate(reader.pages):
                    try:
                        t = page.extract_text() or ""
                    except Exception:
                        t = ""
                    pages_text.append(f"--- Page {i + 1} ---\n{t}")
                text = "\n\n".join(pages_text).strip()
                if text:
                    max_chars = 60_000
                    if len(text) > max_chars:
                        text = text[:max_chars] + f"\n\n[truncated at {max_chars} chars of {len(text)} total]"
                    return {"path": str(p), "content": text}
                return {"error": f"could not extract text from PDF: {p.name} (scanned/image-only PDF?)"}
            except Exception as e:
                return {"error": f"failed to parse PDF {p.name}: {e}"}

        # .docx — extract text via python-docx (binary, not UTF-8 readable)
        if str(p).lower().endswith(".docx"):
            try:
                text = _extract_docx_text(p.read_bytes())
                if text:
                    max_chars = 60_000
                    if len(text) > max_chars:
                        text = text[:max_chars] + f"\n\n[truncated at {max_chars} chars of {len(text)} total]"
                    return {"path": str(p), "content": text}
                return {"error": f"could not extract text from docx: {p.name}"}
            except Exception as e:
                return {"error": f"failed to parse docx {p.name}: {e}"}

        try:
            content = p.read_text(
                encoding="utf-8", errors="replace",
            )[:_MAX_READ_BYTES]

            lines = content.split("\n")
            total_lines = len(lines)
            start = offset - 1
            end = start + limit
            selected = lines[start:end]

            result_text = "\n".join(selected)
            return {
                "path": str(p),
                "content": result_text,
                "lines": len(selected),
                "total_lines": total_lines,
                "offset": offset,
                "truncated": end < total_lines,
            }
        except Exception as e:
            return {"error": f"{type(e).__name__}: {e}"}


# --------------------------------------------------------------------------- #
# WriteFileTool
# --------------------------------------------------------------------------- #


class WriteFileTool(Tool):
    """Write text content to a file (creates or overwrites)."""

    name = "write_file"
    description = (
        "Write text content to a file. Creates parent directories. "
        "Overwrites existing content. Path must be in allowed dirs "
        "(active workspace, user home, Downloads, Desktop, preview). "
        "Relative paths (e.g. analysis.md) resolve into the workspace root. "
        "For 下载目录 use absolute path like /Users/<name>/Downloads/file.md."
    )

    def __init__(self, allowed_dirs: Sequence[str | Path] | None = None) -> None:
        self._allowed_dirs = list(allowed_dirs or [os.getcwd()])

    @property
    def parameters_schema(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Path to the file to write.",
                },
                "content": {
                    "type": "string",
                    "description": "The text content to write.",
                },
            },
            "required": ["path", "content"],
        }

    def __call__(self, **kwargs: Any) -> dict[str, Any]:
        path_str = kwargs.get("path", "")
        content = kwargs.get("content", "")

        if not path_str:
            return {"error": "missing 'path'"}
        # LLM sometimes hallucinates a dict where a string is expected
        # (e.g. {"path": "x.md", ...} leaking out as the path). Bail
        # explicitly with a clear error so the model can retry with a
        # string instead of triggering a confusing 'File name too long'
        # OSError from a 2 KB "path".
        if not isinstance(path_str, str):
            return {"error": f"'path' must be a string, got {type(path_str).__name__}: {str(path_str)[:80]!r}"}
        if len(content) > _MAX_WRITE_BYTES:
            return {"error": f"content too large ({len(content)} > {_MAX_WRITE_BYTES} bytes)"}

        try:
            p = _resolve_in_allowlist(path_str, self._allowed_dirs)
        except PermissionError as e:
            logger.info("write_file denied: %s", e)
            return {"error": str(e)}

        try:
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(content, encoding="utf-8")
            logger.info("write_file: wrote %d bytes to %s", len(content), p)
            return {
                "path": str(p),
                "bytes": len(content),
                "status": "ok",
            }
        except Exception as e:
            return {"error": f"{type(e).__name__}: {e}"}


# --------------------------------------------------------------------------- #
# EditFileTool
# --------------------------------------------------------------------------- #


class EditFileTool(Tool):
    """Find-and-replace within a file.

    Replaces the first occurrence of ``old_text`` with ``new_text``.
    If ``old_text`` is not found, returns an error.
    """

    name = "edit_file"
    description = (
        "Find and replace text within a file. Replaces first match. "
        "Use for targeted edits without rewriting the whole file."
    )

    def __init__(self, allowed_dirs: Sequence[str | Path] | None = None) -> None:
        self._allowed_dirs = list(allowed_dirs or [os.getcwd()])

    @property
    def parameters_schema(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Path to the file to edit.",
                },
                "old_text": {
                    "type": "string",
                    "description": "The exact text to find.",
                },
                "new_text": {
                    "type": "string",
                    "description": "The replacement text.",
                },
            },
            "required": ["path", "old_text", "new_text"],
        }

    def __call__(self, **kwargs: Any) -> dict[str, Any]:
        path_str = kwargs.get("path", "")
        old_text = kwargs.get("old_text", "")
        new_text = kwargs.get("new_text", "")

        if not path_str:
            return {"error": "missing 'path'"}
        if not old_text:
            return {"error": "missing 'old_text'"}

        try:
            p = _resolve_in_allowlist(path_str, self._allowed_dirs)
        except PermissionError as e:
            logger.info("edit_file denied: %s", e)
            return {"error": str(e)}

        if not p.exists():
            return {"error": f"file not found: {p}"}

        try:
            content = p.read_text(encoding="utf-8", errors="replace")

            # Fast path: byte-exact substring (keeps the historical
            # behaviour and cost model).
            if old_text in content:
                new_content = content.replace(old_text, new_text, 1)
                p.write_text(new_content, encoding="utf-8")
                logger.info("edit_file: replaced %d chars in %s", len(old_text), p)
                return {
                    "path": str(p),
                    "status": "ok",
                    "match": "exact",
                    "old_len": len(old_text),
                    "new_len": len(new_text),
                }

            # Fuzzy path (codex apply-patch seek_sequence port): locate
            # old_text as a line sequence with decreasing strictness —
            # exact lines → rstrip → trim → Unicode punctuation
            # normalisation. Rescues model-authored anchors that differ
            # from the file only in trailing whitespace or typographic
            # quotes/dashes instead of failing the whole edit.
            eol = "\r\n" if "\r\n" in content else "\n"
            lines = content.split(eol)
            pattern = old_text.replace("\r\n", "\n").split("\n")
            idx, tier = _seek_sequence(lines, pattern)
            if idx is None:
                return {
                    "error": (
                        f"old_text not found in {p}（精确与模糊匹配均失败）。"
                        "请 read_file 后用文件中的原文重试——注意保留缩进与标点。"
                    ),
                }
            replacement = new_text.replace("\r\n", "\n").split("\n")
            new_lines = lines[:idx] + replacement + lines[idx + len(pattern):]
            new_content = eol.join(new_lines)
            p.write_text(new_content, encoding="utf-8")
            logger.info("edit_file: fuzzy(%s) replaced %d lines in %s",
                        tier, len(pattern), p)
            return {
                "path": str(p),
                "status": "ok",
                "match": tier,
                "old_len": len(old_text),
                "new_len": len(new_text),
            }
        except Exception as e:
            return {"error": f"{type(e).__name__}: {e}"}


# Unicode punctuation normalisation (codex apply-patch): typographic
# dashes/quotes/spaces collapse to their ASCII equivalents so anchors
# authored as plain ASCII still match files containing smart glyphs.
_PUNCT_MAP = str.maketrans({
    "\u2010": "-", "\u2011": "-", "\u2012": "-", "\u2013": "-",
    "\u2014": "-", "\u2015": "-", "\u2212": "-",
    "\u2018": "'", "\u2019": "'", "\u201a": "'", "\u201b": "'",
    "\u201c": '"', "\u201d": '"', "\u201e": '"', "\u201f": '"',
    "\u00a0": " ", "\u2002": " ", "\u2003": " ", "\u2004": " ",
    "\u2005": " ", "\u2006": " ", "\u2007": " ", "\u2008": " ",
    "\u2009": " ", "\u200a": " ", "\u202f": " ", "\u205f": " ",
    "\u3000": " ",
})


def _norm_line(s: str) -> str:
    return s.strip().translate(_PUNCT_MAP)


def _seek_sequence(lines: list[str], pattern: list[str]) -> tuple[int | None, str]:
    """Find `pattern` (a line sequence) in `lines` with decreasing
    strictness. Returns (start_index, tier) or (None, "").

    Tiers, mirroring codex-rs/apply-patch/src/seek_sequence.rs:
      exact → rstrip (ignore trailing whitespace) → trim (both sides)
      → Unicode-punctuation-normalised trim. Empty pattern matches at 0.
    """
    if not pattern:
        return 0, "exact"
    if len(pattern) > len(lines):
        return None, ""
    n = len(lines) - len(pattern)

    for i in range(n + 1):
        if lines[i:i + len(pattern)] == pattern:
            return i, "exact"
    for i in range(n + 1):
        if all(lines[i + j].rstrip() == pattern[j].rstrip()
               for j in range(len(pattern))):
            return i, "rstrip"
    for i in range(n + 1):
        if all(lines[i + j].strip() == pattern[j].strip()
               for j in range(len(pattern))):
            return i, "trim"
    for i in range(n + 1):
        if all(_norm_line(lines[i + j]) == _norm_line(pattern[j])
               for j in range(len(pattern))):
            return i, "unicode"
    return None, ""


# --------------------------------------------------------------------------- #
# WriteXlsxTool
# --------------------------------------------------------------------------- #


class WriteXlsxTool(Tool):
    """Generate a new .xlsx spreadsheet from structured data.

    Lets the agent *produce* spreadsheets, not just read them. The model
    supplies a list of sheets, each with a name and a list of rows (rows are
    lists of cell values). Paths are confined to the allowlist like the
    other file tools.
    """

    name = "write_xlsx"
    description = (
        "Generate a new .xlsx spreadsheet file from structured data. "
        "Provide `sheets`: a list where each item has `name` (sheet name) and "
        "`rows` (a list of rows, each row a list of cell values). "
        "Path must be inside allowed dirs."
    )

    def __init__(self, allowed_dirs: Sequence[str | Path] | None = None) -> None:
        self._allowed_dirs = list(allowed_dirs or [os.getcwd()])

    @property
    def parameters_schema(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Output .xlsx path (e.g. /workspace/report.xlsx).",
                },
                "sheets": {
                    "type": "array",
                    "description": (
                        "Sheets to write. Each item: "
                        "{'name': str, 'rows': list[list[cell]]}."
                    ),
                    "items": {
                        "type": "object",
                        "properties": {
                            "name": {"type": "string"},
                            "rows": {
                                "type": "array",
                                "items": {"type": "array"},
                            },
                        },
                        "required": ["name", "rows"],
                    },
                },
            },
            "required": ["path", "sheets"],
        }

    def __call__(self, **kwargs: Any) -> dict[str, Any]:
        path_str = kwargs.get("path", "")
        sheets = kwargs.get("sheets", [])
        if not path_str:
            return {"error": "missing 'path'"}
        if not isinstance(sheets, list) or not sheets:
            return {"error": "missing or invalid 'sheets' (expected a non-empty list)"}

        try:
            p = _resolve_in_allowlist(path_str, self._allowed_dirs)
        except PermissionError as e:
            logger.info("write_xlsx denied: %s", e)
            return {"error": str(e)}

        try:
            import openpyxl as _xl
            wb = _xl.Workbook()
            for i, sh in enumerate(sheets):
                name = str(sh.get("name", f"Sheet{i + 1}"))[:31]
                rows = sh.get("rows", []) or []
                ws = wb.active if i == 0 else wb.create_sheet(title=name)
                if i == 0:
                    ws.title = name
                for row in rows:
                    if not isinstance(row, list):
                        row = [row]
                    ws.append(["" if c is None else c for c in row])
            wb.save(str(p))
            wb.close()
            return {
                "path": str(p),
                "status": "ok",
                "sheets": len(sheets),
                "bytes": p.stat().st_size,
            }
        except Exception as e:
            return {"error": f"{type(e).__name__}: {e}"}


class WritePptxTool(Tool):
    """Generate a real .pptx deck from a slide outline.

    The model previously tried to build decks through ad-hoc bash +
    python-pptx one-liners, failed (broken shell tool), and then told
    the user the file was written anyway. This tool makes deck
    generation a first-class call: the model supplies a title and a
    slide outline; we render a clean, consistent deck via python-pptx.
    """

    name = "write_pptx"
    description = (
        "Generate a new .pptx presentation file. Provide `title` (deck "
        "title) and `slides`: a list where each item has `title` and "
        "`bullets` (a list of strings; optionally `notes`). Path must be "
        "inside allowed dirs. Renders a clean 16:9 deck with consistent "
        "typography."
    )

    def __init__(self, allowed_dirs: Sequence[str | Path] | None = None) -> None:
        self._allowed_dirs = list(allowed_dirs or [os.getcwd()])

    @property
    def parameters_schema(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Output .pptx path (e.g. /workspace/report.pptx).",
                },
                "title": {"type": "string", "description": "Deck title (cover slide)."},
                "subtitle": {"type": "string", "description": "Optional cover subtitle."},
                "slides": {
                    "type": "array",
                    "description": (
                        "Slide outline. Each item: {'title': str, "
                        "'bullets': list[str], 'notes': str (optional)}."
                    ),
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "bullets": {"type": "array", "items": {"type": "string"}},
                            "notes": {"type": "string"},
                        },
                        "required": ["title"],
                    },
                },
            },
            "required": ["path", "title", "slides"],
        }

    def __call__(self, **kwargs: Any) -> dict[str, Any]:
        path_str = kwargs.get("path", "")
        title = str(kwargs.get("title", "") or "").strip()
        slides = kwargs.get("slides", [])
        if not path_str:
            return {"error": "missing 'path'"}
        if not title:
            return {"error": "missing 'title'"}
        if not isinstance(slides, list) or not slides:
            return {"error": "missing or invalid 'slides' (expected a non-empty list)"}

        try:
            p = _resolve_in_allowlist(path_str, self._allowed_dirs)
        except PermissionError as e:
            logger.info("write_pptx denied: %s", e)
            return {"error": str(e)}

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            prs = Presentation()
            prs.slide_width = Inches(13.333)   # 16:9
            prs.slide_height = Inches(7.5)

            ACCENT = "1A1A1E"      # ink
            ACCENT_SOFT = "5B5B66"
            # ── Cover ──
            cover = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            bar = cover.shapes.add_shape(1, Inches(0.9), Inches(2.35), Inches(0.14), Inches(1.5))
            bar.fill.solid(); bar.fill.fore_color.rgb = __import__("pptx.dml.color", fromlist=["RGBColor"]).RGBColor(0x1A, 0x1A, 0x1E)
            bar.line.fill.background()
            box = cover.shapes.add_textbox(Inches(1.25), Inches(2.3), Inches(11.0), Inches(1.3))
            para = box.text_frame.paragraphs[0]
            run = para.add_run(); run.text = title
            run.font.size = Pt(44); run.font.bold = True
            run.font.color.rgb = __import__("pptx.dml.color", fromlist=["RGBColor"]).RGBColor(0x1A, 0x1A, 0x1E)
            sub = str(kwargs.get("subtitle", "") or "").strip()
            if sub:
                sbox = cover.shapes.add_textbox(Inches(1.27), Inches(3.55), Inches(11.0), Inches(0.7))
                sp = sbox.text_frame.paragraphs[0]
                srun = sp.add_run(); srun.text = sub
                srun.font.size = Pt(18)
                srun.font.color.rgb = __import__("pptx.dml.color", fromlist=["RGBColor"]).RGBColor(0x5B, 0x5B, 0x66)

            # ── Content slides ──
            for i, sl in enumerate(slides, start=1):
                if not isinstance(sl, dict):
                    sl = {"title": str(sl)}
                stitle = str(sl.get("title", f"Slide {i}"))
                bullets = sl.get("bullets", []) or []
                notes = str(sl.get("notes", "") or "")
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                tbox = slide.shapes.add_textbox(Inches(0.9), Inches(0.55), Inches(11.5), Inches(1.0))
                tp = tbox.text_frame.paragraphs[0]
                trun = tp.add_run(); trun.text = stitle
                trun.font.size = Pt(30); trun.font.bold = True
                trun.font.color.rgb = __import__("pptx.dml.color", fromlist=["RGBColor"]).RGBColor(0x1A, 0x1A, 0x1E)
                rule = slide.shapes.add_shape(1, Inches(0.93), Inches(1.55), Inches(12.0), Pt(2))
                rule.fill.solid(); rule.fill.fore_color.rgb = __import__("pptx.dml.color", fromlist=["RGBColor"]).RGBColor(0xE4, 0xE4, 0xE7)
                rule.line.fill.background()
                if bullets:
                    bbox = slide.shapes.add_textbox(Inches(0.95), Inches(1.9), Inches(11.4), Inches(5.0))
                    tf = bbox.text_frame
                    tf.word_wrap = True
                    for bi, b in enumerate(bullets):
                        bp = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
                        br = bp.add_run()
                        br.text = f"•  {b}" if not str(b).startswith(("•", "-", "数字")) else str(b)
                        br.font.size = Pt(17)
                        br.font.color.rgb = __import__("pptx.dml.color", fromlist=["RGBColor"]).RGBColor(0x27, 0x27, 0x2A)
                        bp.space_after = Pt(10)
                if notes:
                    slide.notes_slide.notes_text_frame.text = notes

            # footer page numbers
            for i, slide in enumerate(prs.slides):
                if i == 0:
                    continue
                nbox = slide.shapes.add_textbox(Inches(12.3), Inches(7.02), Inches(0.8), Inches(0.35))
                np_ = nbox.text_frame.paragraphs[0]
                nrun = np_.add_run(); nrun.text = str(i)
                nrun.font.size = Pt(11)
                nrun.font.color.rgb = __import__("pptx.dml.color", fromlist=["RGBColor"]).RGBColor(0x8A, 0x8A, 0x93)

            prs.save(str(p))
            return {
                "path": str(p),
                "status": "ok",
                "slides": len(slides) + 1,
                "bytes": p.stat().st_size,
            }
        except Exception as e:
            return {"error": f"{type(e).__name__}: {e}"}


class ReadOfficeTool(Tool):
    """Read an office document (xlsx/pptx/docx) as structured text.

    Binary office formats were a one-way door: the model could WRITE
    them (write_xlsx/write_pptx) but read_file returned mojibake, so it
    could never verify its own output or iterate. This tool converts
    each format to markdown-ish text the model can actually reason
    about — the read half of the office loop (Claude-Skills style).
    """

    name = "read_office"
    description = (
        "Read an office document and return its content as text. "
        "Supports .xlsx (sheets as markdown tables), .pptx (slide "
        "titles + bullets), .docx (paragraphs + tables). Use this to "
        "inspect a document before editing it, or to verify a file you "
        "just generated."
    )

    def __init__(self, allowed_dirs: Sequence[str | Path] | None = None) -> None:
        self._allowed_dirs = [Path(d).expanduser().resolve() for d in (allowed_dirs or [os.getcwd()])]

    @property
    def parameters_schema(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Path to the .xlsx/.pptx/.docx file.",
                },
                "max_chars": {
                    "type": "integer",
                    "description": "Output cap (default 6000). Raise for big files.",
                },
            },
            "required": ["path"],
        }

    # ── converters ──
    @staticmethod
    def _read_xlsx(path: Path) -> str:
        import openpyxl as _xl
        wb = _xl.load_workbook(str(path), data_only=True, read_only=True)
        out = []
        for ws in wb.worksheets:
            out.append(f"## Sheet: {ws.title} ({ws.max_row} rows x {ws.max_column} cols)")
            for row in ws.iter_rows(max_row=200, values_only=True):
                cells = ["" if c is None else str(c) for c in row]
                if any(c.strip() for c in cells):
                    out.append("| " + " | ".join(cells) + " |")
            out.append("")
        wb.close()
        return "\n".join(out)

    @staticmethod
    def _read_pptx(path: Path) -> str:
        from pptx import Presentation
        prs = Presentation(str(path))
        out = [f"Presentation: {len(prs.slides.__iter__.__self__._sldIdLst)} slides, "
               f"{prs.slide_width.inches:.1f}x{prs.slide_height.inches:.1f} in"]
        for i, slide in enumerate(prs.slides, start=1):
            out.append(f"## Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    out.append(shape.text_frame.text.strip())
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                out.append(f"[notes] {slide.notes_slide.notes_text_frame.text.strip()[:300]}")
        return "\n".join(out)

    @staticmethod
    def _read_docx(path: Path) -> str:
        import docx as _dx
        doc = _dx.Document(str(path))
        out = []
        for para in doc.paragraphs:
            if para.text.strip():
                style = para.style.name if para.style else ""
                prefix = "#" * min(4, max(1, len(style) - len("Heading") + 1)) if style.startswith("Heading") else ""
                out.append(f"{prefix + ' ' if prefix else ''}{para.text.strip()}")
        for t_idx, table in enumerate(doc.tables, start=1):
            out.append(f"## Table {t_idx}")
            for row in table.rows[:100]:
                cells = [c.text.strip() for c in row.cells]
                out.append("| " + " | ".join(cells) + " |")
        return "\n".join(out)

    def __call__(self, **kwargs: Any) -> dict[str, Any]:
        path_str = kwargs.get("path", "")
        max_chars = int(kwargs.get("max_chars", 6000) or 6000)
        if not path_str:
            return {"error": "missing 'path'"}
        try:
            p = _resolve_in_allowlist(path_str, self._allowed_dirs)
        except PermissionError as e:
            return {"error": str(e)}
        if not p.exists():
            return {"error": f"file not found: {p}"}
        suffix = p.suffix.lower()
        try:
            if suffix == ".xlsx":
                text = self._read_xlsx(p)
            elif suffix == ".pptx":
                text = self._read_pptx(p)
            elif suffix == ".docx":
                text = self._read_docx(p)
            else:
                return {"error": f"unsupported type '{suffix}' — read_office handles .xlsx/.pptx/.docx"}
        except Exception as e:
            return {"error": f"{type(e).__name__}: {e}"}
        truncated = len(text) > max_chars
        return {
            "path": str(p),
            "type": suffix.lstrip("."),
            "content": text[:max_chars],
            "truncated": truncated,
            "total_chars": len(text),
        }


__all__ = [
    "ReadFileTool",
    "WriteFileTool",
    "EditFileTool",
    "WriteXlsxTool",
    "WritePptxTool",
    "ReadOfficeTool",
]
